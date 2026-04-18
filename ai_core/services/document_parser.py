"""Document text extraction and chunking utilities.

Supports PDF, DOCX, and PPTX file formats. Extracted text is split into
overlapping chunks suitable for embedding and semantic search.
"""

import io
import logging
from pathlib import Path
from typing import Union

logger = logging.getLogger("ai_core")


def extract_text(file_input: Union[str, io.IOBase]) -> str:
    """Extract plain text content from a document file.

    Args:
        file_input: Either a file path (str) or file-like object (e.g., UploadedFile, BytesIO).

    Returns:
        Extracted text as a single string. Returns empty string
        if the file format is unsupported or extraction fails.
    """
    # Determine file extension
    if isinstance(file_input, str):
        ext = Path(file_input).suffix.lower()
    else:
        # File object - try to get name attribute, fallback to stream detection
        name = getattr(file_input, 'name', '')
        ext = Path(name).suffix.lower() if name else ''

    try:
        if ext == ".pdf":
            return _extract_pdf(file_input)
        elif ext in (".doc", ".docx"):
            return _extract_docx(file_input)
        elif ext == ".pptx":
            return _extract_pptx(file_input)
        else:
            logger.warning("Unsupported file format for text extraction: %s", ext)
            return ""
    except Exception as e:
        logger.error("Failed to extract text from %s: %s", file_input, e)
        return ""


def chunk_text(
    text: str, chunk_size: int = 4000, overlap: int = 200
) -> list[str]:
    """Split text into overlapping chunks for embedding.

    Args:
        text: The full text to split.
        chunk_size: Maximum number of characters per chunk.
        overlap: Number of overlapping characters between consecutive chunks.

    Returns:
        A list of text chunks. Returns empty list if input text is blank.
    """
    if not text or not text.strip():
        return []

    chunks: list[str] = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = min(start + chunk_size, text_length)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap

    logger.info("Split text into %d chunks (size=%d, overlap=%d)", len(chunks), chunk_size, overlap)
    return chunks


# ---------------------------------------------------------------------------
# Private extraction helpers
# ---------------------------------------------------------------------------


def _extract_pdf(file_input: Union[str, io.IOBase]) -> str:
    """Extract text from a PDF file using PyPDF2."""
    from PyPDF2 import PdfReader

    reader = PdfReader(file_input)
    pages: list[str] = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            pages.append(page_text)

    text = "\n".join(pages)
    logger.info("Extracted %d characters from PDF (%d pages)", len(text), len(reader.pages))
    return text


def _extract_docx(file_input: Union[str, io.IOBase]) -> str:
    """Extract text from a DOCX file using python-docx."""
    from docx import Document

    doc = Document(file_input)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    text = "\n".join(paragraphs)
    logger.info("Extracted %d characters from DOCX (%d paragraphs)", len(text), len(paragraphs))
    return text


def _extract_pptx(file_input: Union[str, io.IOBase]) -> str:
    """Extract text from a PPTX file using python-pptx."""
    from pptx import Presentation

    prs = Presentation(file_input)
    slides_text: list[str] = []

    for slide in prs.slides:
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        slide_parts.append(paragraph_text)
        if slide_parts:
            slides_text.append("\n".join(slide_parts))

    text = "\n\n".join(slides_text)
    logger.info("Extracted %d characters from PPTX (%d slides)", len(text), len(prs.slides))
    return text
